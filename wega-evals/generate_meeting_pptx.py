"""Generate meeting PowerPoint: Wega-Evals Review & Assessment Plan."""

from __future__ import annotations
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colour palette ──────────────────────────────────────────────────
BRAND_DARK  = RGBColor(0x1B, 0x2A, 0x4A)
BRAND_MID   = RGBColor(0x2E, 0x86, 0xAB)
BRAND_LIGHT = RGBColor(0xE8, 0xF1, 0xF5)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
GREY        = RGBColor(0x7F, 0x8C, 0x8D)
GREEN       = RGBColor(0x27, 0xAE, 0x60)
RED         = RGBColor(0xE7, 0x4C, 0x3C)
ORANGE      = RGBColor(0xF3, 0x9C, 0x12)
LIGHT_GREY  = RGBColor(0xEC, 0xF0, 0xF1)

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def _set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_bullet_slide(prs, title_text, bullets, sub_title=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, WHITE)
    # Top bar
    _add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), BRAND_DARK)
    _add_textbox(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
                 title_text, font_size=28, color=WHITE, bold=True)

    if sub_title:
        _add_textbox(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
                     sub_title, font_size=16, color=BRAND_MID, bold=True)

    # Bullets
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8 if sub_title else 1.4),
                                     Inches(11.5), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK
        p.font.name = "Calibri"
        p.space_after = Pt(8)
        p.level = 0
        # Bullet character
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '▸'})
        # Remove existing bullets
        for old in pPr.findall(qn('a:buChar')):
            pPr.remove(old)
        for old in pPr.findall(qn('a:buNone')):
            pPr.remove(old)
        pPr.append(buChar)

    return slide


def _add_table_slide(prs, title_text, headers, rows, col_widths=None, sub_title=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    # Top bar
    _add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), BRAND_DARK)
    _add_textbox(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
                 title_text, font_size=28, color=WHITE, bold=True)

    if sub_title:
        _add_textbox(slide, Inches(0.6), Inches(1.15), Inches(12), Inches(0.4),
                     sub_title, font_size=14, color=BRAND_MID, bold=False)

    num_rows = len(rows) + 1
    num_cols = len(headers)
    top = Inches(1.7) if sub_title else Inches(1.4)
    table_width = Inches(12)
    table_height = Inches(0.4) * num_rows

    table_shape = slide.shapes.add_table(num_rows, num_cols,
                                          Inches(0.6), top,
                                          table_width, table_height)
    table = table_shape.table

    # Column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    # Header
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BRAND_DARK
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.name = "Calibri"

    # Rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            if r_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GREY
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = BLACK
                p.font.name = "Calibri"

    return slide


def generate_pptx():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 1 — Title
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BRAND_DARK)

    _add_textbox(slide, Inches(1), Inches(2), Inches(11), Inches(1.2),
                 "Wega-Evals", font_size=48, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
                 "Evaluation Framework Review & Agent Assessment Plan",
                 font_size=24, color=BRAND_MID, bold=False, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
                 datetime.now().strftime("%B %d, %Y"),
                 font_size=16, color=GREY, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 2 — Agenda
    # ══════════════════════════════════════════════════════════════
    _add_bullet_slide(prs, "Agenda", [
        "1. Review Wega-Evals evaluation framework and architecture",
        "2. Discuss uniqueness versus off-the-shelf model usage",
        "3. Evaluate agent assessment plan and golden dataset creation",
        "4. Provide feedback on evaluation approach and deployment structure",
    ])

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 3 — What is Wega-Evals
    # ══════════════════════════════════════════════════════════════
    _add_bullet_slide(prs, "1. What is Wega-Evals?", [
        "Profile-driven evaluation framework for 11 SDLC AI agents",
        "Combines LLM-as-Judge (Gemini 2.5) + Programmatic metrics",
        "Scores each agent across 15–20 dimensions per evaluation",
        "Evaluates real production traces from Langfuse — not synthetic benchmarks",
        "Plugin architecture — add new agent evaluation in < 1 hour",
        "Built-in safety: hallucination, toxicity, PII, policy compliance auto-applied",
    ])

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 4 — Architecture Layers
    # ══════════════════════════════════════════════════════════════
    _add_table_slide(prs, "Architecture — 4-Layer Design", 
        ["Layer", "Components", "Role"],
        [
            ["Presentation", "CLI, Streamlit Dashboard, DOCX Reports", "User interaction & visualization"],
            ["Orchestration", "EvalRunner (ThreadPoolExecutor)", "Parallel processing, score aggregation"],
            ["Evaluation Engine", "LLM Judge + Programmatic Evals", "Dual scoring: AI judgment + deterministic metrics"],
            ["Data & Integration", "DatasetManager, Langfuse, Vertex AI", "Trace storage, datasets, model hosting"],
        ],
        col_widths=[2.5, 4.5, 5.0],
    )

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 5 — How It Works
    # ══════════════════════════════════════════════════════════════
    _add_bullet_slide(prs, "How Evaluation Works", [
        "Step 1: Production traces captured in Langfuse from real agent usage",
        "Step 2: SEED — traces converted into evaluation dataset items",
        "Step 3: EVALUATE — each item scored across ~17 dimensions",
        "        ▸ LLM Judge: Gemini 2.5 scores quality, faithfulness, safety (0.0–1.0)",
        "        ▸ Programmatic: latency, cost, token efficiency, structural validity",
        "Step 4: REPORT — scores pushed to Langfuse, dashboard, DOCX exports",
        "Two-level parallelism: 4 items × 8 dimensions = up to 32 concurrent evals",
    ])

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 6 — Evaluation Dimensions
    # ══════════════════════════════════════════════════════════════
    _add_table_slide(prs, "26+ Evaluation Dimensions (5 Categories)",
        ["Category", "Dimensions", "Evaluator"],
        [
            ["Output Quality", "faithfulness, coherence, conciseness, consistency", "LLM Judge"],
            ["Safety & Trust", "hallucination, toxicity, data_privacy_compliance", "LLM Judge"],
            ["Compliance", "policy_compliance", "LLM Judge"],
            ["Performance", "latency, cost_efficiency, token_efficiency, structural_validity", "Programmatic"],
            ["Domain-Specific", "script_correctness, review_relevance, intent_routing, etc.", "LLM Judge"],
        ],
        col_widths=[2.5, 6.5, 3.0],
        sub_title="8 universal + 4 programmatic + 5–8 domain-specific per agent"
    )

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 7 — Technology Stack
    # ══════════════════════════════════════════════════════════════
    _add_table_slide(prs, "Technology Stack",
        ["Component", "Technology", "Purpose"],
        [
            ["Runtime", "Python ≥3.11", "Framework language"],
            ["LLM Judge", "Google Gemini 2.5 via Vertex AI", "Structured evaluation scoring"],
            ["Observability", "Langfuse 4.x (self-hosted K8s)", "Trace storage, dataset mgmt, scores"],
            ["Data Models", "Pydantic 2.x", "Type-safe data contracts"],
            ["Dashboard", "Streamlit", "Real-time evaluation visualization"],
            ["Reports", "python-docx + matplotlib", "Professional DOCX with charts"],
            ["Config", "python-dotenv + dataclass", "12-factor environment-driven config"],
        ],
        col_widths=[2.5, 4.5, 5.0],
    )

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 8 — Section Divider: Uniqueness
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BRAND_MID)
    _add_textbox(slide, Inches(1), Inches(2.8), Inches(11), Inches(1),
                 "2. Uniqueness vs. Off-the-Shelf",
                 font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 9 — Composition Breakdown
    # ══════════════════════════════════════════════════════════════
    _add_table_slide(prs, "Composition Breakdown",
        ["Category", "%", "What It Covers"],
        [
            ["Custom / Unique", "~40%", "Agent-specific profiles, domain dimensions, custom extractors, tailored prompts"],
            ["Standard Patterns", "~40%", "LLM-as-Judge approach, Pydantic models, CLI scaffolding, dashboard patterns"],
            ["Glue / Integration", "~20%", "Langfuse connectivity, SSL adapter, Vertex AI wiring, report generation"],
        ],
        col_widths=[3, 1.5, 7.5],
    )

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 10 — What Makes It Unique
    # ══════════════════════════════════════════════════════════════
    _add_bullet_slide(prs, "What Makes Wega-Evals Unique", [
        "SDLC-domain evaluation taxonomy — 26+ dimensions for code review, test gen, BRD, orchestration",
        "Profile-driven architecture — one framework evaluates 11 different agents, zero core code changes",
        "Production trace evaluation — real Langfuse traces, not synthetic benchmarks",
        "Auto-merge dimension system — __init_subclass__ ensures every agent gets safety/quality baselines",
        "Corporate environment ready — self-signed cert chain builder baked in",
    ], sub_title="Cannot be purchased off-the-shelf")

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 11 — Market Comparison
    # ══════════════════════════════════════════════════════════════
    _add_table_slide(prs, "Market Comparison — Off-the-Shelf Options",
        ["Capability", "Available Tools", "Why Not Sufficient"],
        [
            ["LLM-as-Judge", "Ragas, DeepEval, LangSmith, Phoenix", "Focus on single LLM calls, not multi-step SDLC agents"],
            ["Trace Storage", "Langfuse, LangSmith, Arize Phoenix", "We already use Langfuse — this is the data source, not the evaluator"],
            ["Prompt Testing", "Promptfoo, Humanloop", "Tests prompts, not end-to-end agent pipelines"],
            ["Agent Eval", "None purpose-built for SDLC", "No tool evaluates 11 heterogeneous SDLC agents"],
        ],
        col_widths=[2.5, 4.5, 5.0],
    )

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 12 — Key Message
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.1), BRAND_DARK)
    _add_textbox(slide, Inches(0.6), Inches(0.2), Inches(12), Inches(0.7),
                 "Key Takeaway", font_size=28, color=WHITE, bold=True)

    # Quote box
    _add_shape_rect(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(2.5), BRAND_LIGHT)
    _add_textbox(slide, Inches(2), Inches(2.5), Inches(9), Inches(2),
                 '"We use industry-standard patterns (LLM-as-Judge, Langfuse, Pydantic) '
                 'but the evaluation logic — the dimensions, prompts, and extraction rules — '
                 'is domain-specific and cannot be bought.\n\n'
                 'The framework is the reusable shell; the profiles are the intellectual property."',
                 font_size=20, color=BRAND_DARK, bold=False, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 13 — Section Divider: Assessment Plan
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BRAND_MID)
    _add_textbox(slide, Inches(1), Inches(2.8), Inches(11), Inches(1),
                 "3. Agent Assessment Plan & Golden Dataset",
                 font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 14 — Agent Readiness
    # ══════════════════════════════════════════════════════════════
    _add_table_slide(prs, "Agent Readiness — Current Status",
        ["Priority", "Agent", "Traces", "Status", "Next Action"],
        [
            ["✅ Done", "testcase_to_scripts", "171", "Evaluated (100 items)", "Review results"],
            ["✅ Done", "testcases_to_testdata", "96", "Evaluated (85 items)", "Review results"],
            ["1", "code_assistant", "193", "Not seeded", "Seed & evaluate"],
            ["2", "brd_summary", "29", "Not seeded", "Seed & evaluate"],
            ["3", "brd", "23", "Not seeded", "Seed & evaluate"],
            ["4", "user_manual", "17", "Needs re-seed", "Re-seed & evaluate"],
            ["5", "sdlc_orchestrator", "1,309", "Needs re-seed", "Seed & evaluate (long run)"],
            ["6", "cara", "4", "Needs re-seed", "Small sample"],
            ["—", "user_story", "0", "No traces", "Needs production usage"],
            ["—", "userstory_to_testcases", "0", "No traces", "Needs production usage"],
            ["—", "userstory_validator", "0", "No traces", "Needs production usage"],
        ],
        col_widths=[1.3, 3.0, 1.2, 3.0, 3.5],
    )

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 15 — Completed Evaluation Results
    # ══════════════════════════════════════════════════════════════
    _add_table_slide(prs, "Completed Evaluation Results",
        ["Agent", "Items", "Dimensions", "Overall Score", "Status", "Top Performer", "Needs Work"],
        [
            ["testcase_to_scripts", "100", "17", "0.59", "NEEDS IMPROVEMENT",
             "structural_validity (1.0), toxicity (1.0)", "hallucination (0.18), faithfulness (0.32)"],
            ["testcases_to_testdata", "85", "15", "0.47", "NEEDS IMPROVEMENT",
             "consistency (0.97), toxicity (1.0), latency (1.0)", "boundary_values (0.06), structural_validity (0.06)"],
        ],
        col_widths=[2.2, 0.8, 1.2, 1.2, 2, 2.8, 2.8],
        sub_title="Both agents show strong safety scores but need quality improvements"
    )

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 16 — Golden Dataset Strategy
    # ══════════════════════════════════════════════════════════════
    _add_table_slide(prs, "Golden Dataset Creation — Phased Approach",
        ["Phase", "Action", "Effort", "Outcome"],
        [
            ["Phase 1 (Now)", "Trace-seeded datasets — LLM judges output quality without reference", "Automated", "Baseline quality scores"],
            ["Phase 2", "SMEs review 10–20 outputs per top agent, mark gold-standard answers", "2–3 days/agent", "Human-verified references"],
            ["Phase 3", "Enable reference-comparison scoring — output vs. golden reference", "1 sprint", "Accuracy measurement"],
            ["Phase 4", "Regression suite — re-evaluate on golden dataset after agent updates", "Ongoing", "Quality gates"],
        ],
        col_widths=[2, 5, 2, 3],
    )

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 17 — Zero-Trace Agents
    # ══════════════════════════════════════════════════════════════
    _add_bullet_slide(prs, "Zero-Trace Agents — Action Needed", [
        "3 agents have NO production traces: user_story, userstory_to_testcases, userstory_validator",
        "",
        "Option A: Run them in production/staging to generate Langfuse traces, then seed",
        "Option B: Manually create synthetic dataset items using CLI (wega-eval add-item)",
        "Option C: Targeted test sessions — QA uses agents on 10–15 real requirements",
        "",
        "Recommendation: Option C — generates realistic traces with minimal production risk",
    ], sub_title="No traces = nothing to evaluate")

    # ══════════════════════════════════════════════════════════════
    #  Section Divider: Deep Dive — Manual Work Per Agent
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BRAND_MID)
    _add_textbox(slide, Inches(1), Inches(2.2), Inches(11), Inches(1),
                 "Deep Dive: What Is Built Manually Per Agent",
                 font_size=38, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
                 "Dimensions  •  Prompts  •  Extraction Rules  •  Template Variables  •  Programmatic Checks",
                 font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════
    #  Dimension Layers
    # ══════════════════════════════════════════════════════════════
    _add_table_slide(prs, "1. Dimensions — 3 Layers",
        ["Layer", "Source", "Count", "Examples", "Effort"],
        [
            ["Universal\n(auto-applied)", "common_prompts.py\n(__init_subclass__)", "8", "hallucination, faithfulness, coherence,\nconsistency, toxicity, data_privacy, policy", "Written once,\nnever again"],
            ["Programmatic\n(auto-applied)", "AgentProfile base class", "4", "structural_validity, latency,\ncost_efficiency, token_efficiency", "Written once,\nnever again"],
            ["Domain-Specific\n(manual per agent)", "profiles/<agent>.py", "5–8", "CARA: review_relevance, finding_accuracy,\nseverity_calibration, false_positive_rate", "~10 min\nper agent"],
        ],
        col_widths=[2.2, 2.5, 1, 4.3, 2],
        sub_title="Each agent gets 15–20 total dimensions. Domain expert decides what matters for that specific agent."
    )

    # ══════════════════════════════════════════════════════════════
    #  Judge Prompts
    # ══════════════════════════════════════════════════════════════
    _add_bullet_slide(prs, "2. Judge Prompts — Hand-Crafted Evaluation Rubrics", [
        "Each dimension maps to a detailed prompt template sent to Gemini 2.5",
        "Prompt defines: evaluation criteria (bullet points), scoring rubric (0.0–1.0), JSON output format",
        "Uses template variables: {input_text}, {output_text} + agent-specific like {diff}, {framework_type}",
        "",
        "Example — CARA review_relevance prompt structure:",
        '   ▸ Role: "You are an expert code review evaluator"',
        "   ▸ Context: ## Code Diff {diff}  +  ## AI Review Output {output_text}",
        "   ▸ Criteria: 4 specific questions (findings in diff? correct line numbers? etc.)",
        "   ▸ Rubric: 1.0 = all relevant, 0.5 = mixed, 0.0 = mostly hallucinated",
        '   ▸ Output: JSON { "score": float, "reasoning": str, "hallucinated_findings": [...] }',
        "",
        "Effort: 2–4 hours per agent (5–8 prompts × 15–30 min each)",
    ], sub_title="This is the core intellectual property — domain expertise encoded as evaluation rubrics")

    # ══════════════════════════════════════════════════════════════
    #  Extraction Rules
    # ══════════════════════════════════════════════════════════════
    _add_bullet_slide(prs, "3. Extraction Rules — Finding Data in Langfuse Traces", [
        "Each agent stores data differently in Langfuse traces. Extractors tell the framework WHERE to look.",
        "",
        "Default (base class): trace.input → input, last GENERATION observation → output",
        "",
        "Custom overrides needed when agents have non-standard trace structures:",
        '   ▸ CARA: diff is inside a GENERATION observation\'s input, not trace.input',
        '   ▸ testcase_to_scripts: config (framework, language) is in a SPAN named "convert_test_cases"',
        '   ▸ testcase_to_scripts: scripts spread across 3 observation types (generate-script,',
        '     generate-feature-file, generate-step-definition)',
        "",
        "Requires: inspecting real Langfuse traces to understand the structure",
        "Effort: 1–2 hours per agent (only needed for non-standard agents)",
    ], sub_title="3 methods: extract_input(), extract_output(), extract_primary_input_text()")

    # ══════════════════════════════════════════════════════════════
    #  Template Variables & Programmatic Checks
    # ══════════════════════════════════════════════════════════════
    _add_table_slide(prs, "4. Template Variables & 5. Programmatic Checks",
        ["Agent", "Template Variables\n(format_judge_context)", "Custom Programmatic Checks", "Purpose"],
        [
            ["CARA", "{diff}, {findings_json}", "severity_distribution,\nline_reference_accuracy", "Judge sees actual diff;\ncheck severity spread & valid line refs"],
            ["testcase_to_scripts", "{framework_type}, {language}", "—", "Judge evaluates against\ncorrect framework (Playwright vs Selenium)"],
            ["user_story", "—", "INVEST compliance check", "Programmatic validation of\nstory format rules"],
            ["Others", "— (use defaults)", "— (base class only)", "Standard latency / cost /\nstructure checks only"],
        ],
        col_widths=[2.2, 3.3, 3, 3.5],
        sub_title="Template variables inject agent-specific context into judge prompts. Programmatic checks run without LLM."
    )

    # ══════════════════════════════════════════════════════════════
    #  Effort Summary Per Agent
    # ══════════════════════════════════════════════════════════════
    _add_table_slide(prs, "Effort Summary — Manual Work Per Agent Profile",
        ["Component", "What's Done", "Typical Effort", "Skill Needed"],
        [
            ["Dimensions list", "Choose 5–8 quality aspects specific to the agent's domain", "10 min", "Domain knowledge"],
            ["Judge prompts", "Write detailed evaluation rubrics with criteria + scoring guide", "2–4 hours", "Domain + prompt engineering"],
            ["Extraction rules", "Override extract_input/output for non-standard trace structures", "1–2 hours", "Langfuse trace inspection"],
            ["Template variables", "Provide extra context like {diff}, {framework_type} for prompts", "15 min", "Domain knowledge"],
            ["Programmatic checks", "Add rule-based validators (severity distribution, line refs, etc.)", "1–2 hours", "Engineering"],
            ["Output schema", "Define required fields, status values, severity ranges", "10 min", "Agent API knowledge"],
        ],
        col_widths=[2.5, 4.5, 1.8, 3.2],
        sub_title="Total: ~4–8 hours per new agent. Framework handles everything else (parallelism, Langfuse, LLM calls, scoring, reporting)."
    )

    # ══════════════════════════════════════════════════════════════
    #  What's Automatic vs Manual
    # ══════════════════════════════════════════════════════════════
    _add_table_slide(prs, "Automatic vs. Manual — Clear Separation",
        ["Capability", "Automatic (Framework)", "Manual (Per Agent Profile)"],
        [
            ["Dimensions", "8 universal + 4 programmatic auto-merged", "5–8 domain-specific chosen by expert"],
            ["Prompts", "11 common prompts auto-merged", "5–8 custom rubrics per agent"],
            ["Extraction", "Default trace.input / last GENERATION", "Custom overrides for non-standard traces"],
            ["Scoring", "Gemini API calls, JSON parsing, score extraction", "—"],
            ["Parallelism", "ThreadPoolExecutor (4 items × 8 dims)", "—"],
            ["Langfuse Push", "Auto push per-dimension + overall scores", "—"],
            ["Reporting", "Dashboard + DOCX generation", "—"],
            ["SSL / Proxy", "CA bundle auto-build for corporate envs", "—"],
        ],
        col_widths=[2.2, 5.3, 4.5],
        sub_title="Adding a new agent = one Python file. No changes to framework core."
    )

    # ══════════════════════════════════════════════════════════════
    #  Section Divider: Feedback
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BRAND_MID)
    _add_textbox(slide, Inches(1), Inches(2.8), Inches(11), Inches(1),
                 "4. Evaluation Approach & Deployment",
                 font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 19 — Strengths
    # ══════════════════════════════════════════════════════════════
    _add_table_slide(prs, "Current Strengths",
        ["Strength", "Detail"],
        [
            ["Dual Evaluation", "LLM Judge (subjective quality) + Programmatic (objective metrics)"],
            ["Comprehensive Coverage", "8 universal + 4 programmatic + 5–8 domain dims per agent"],
            ["Safety Built-in", "Hallucination, toxicity, PII, policy compliance auto-applied to ALL agents"],
            ["Reproducible", "Temperature=0.0, deterministic programmatic evals, full traceability"],
            ["Scalable", "Plugin architecture — new agent evaluation in <1 hour, no core changes"],
            ["Corporate-Ready", "SSL/proxy handling, self-hosted Langfuse, Vertex AI integration"],
        ],
        col_widths=[3, 9],
    )

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 20 — Gaps & Improvements
    # ══════════════════════════════════════════════════════════════
    _add_table_slide(prs, "Known Gaps & Proposed Improvements",
        ["Gap", "Impact", "Proposed Fix"],
        [
            ["Layer 1 only (single-agent)", "No handoff evaluation between agents", "Build Layer 2 (transition) & Layer 3 (end-to-end)"],
            ["No weighted scoring", "All dimensions contribute equally", "Per-dimension weights (hallucination=3×, conciseness=0.5×)"],
            ["Single judge model", "Potential Gemini bias", "Multi-judge consensus (Gemini + GPT-4 + Claude)"],
            ["No prompt versioning", "Cannot A/B test template changes", "Langfuse prompt management integration"],
            ["No regression tracking", "Cannot compare Agent v1 vs v2", "Comparative evaluation mode on same inputs"],
        ],
        col_widths=[3, 3.5, 5.5],
    )

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 21 — Deployment Evolution
    # ══════════════════════════════════════════════════════════════
    _add_table_slide(prs, "Deployment Structure — Evolution Plan",
        ["Stage", "Structure", "Benefit"],
        [
            ["Current", "Developer runs CLI locally (wega-eval run)", "Quick iteration, debugging, flexibility"],
            ["Next", "CI/CD integration — nightly or on agent PR merge", "Automated regression detection"],
            ["Future", "Dedicated eval service with API + webhook triggers", "Zero-touch continuous evaluation"],
        ],
        col_widths=[2, 5, 5],
    )

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 22 — Discussion Questions
    # ══════════════════════════════════════════════════════════════
    _add_bullet_slide(prs, "Discussion & Feedback", [
        "1. Should we prioritize Layer 2 (transition eval) or golden dataset creation first?",
        "2. Which agents need the most urgent evaluation — by business impact or trace volume?",
        "3. Is Gemini Flash sufficient as judge, or invest in multi-judge consensus?",
        "4. Do we want CI/CD evaluation integration now, or keep manual runs?",
        "5. What's the plan to generate traces for the 3 zero-trace agents?",
        "6. Should we introduce weighted scoring for critical dimensions?",
    ], sub_title="Questions for the team")

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 23 — Summary / Next Steps
    # ══════════════════════════════════════════════════════════════
    _add_bullet_slide(prs, "Summary & Next Steps", [
        "✅ Framework is operational — 2 agents fully evaluated, 6 ready to evaluate",
        "✅ Architecture is extensible — plugin model with auto-merged safety dimensions",
        "✅ Unique value — domain-specific evaluation taxonomy cannot be bought off-the-shelf",
        "",
        "Next Steps:",
        "   ▸ Seed & evaluate code_assistant (193 traces) — highest priority",
        "   ▸ Seed & evaluate brd_summary (29) and brd (23) — quick wins",
        "   ▸ Generate traces for 3 zero-trace agents via targeted QA sessions",
        "   ▸ Begin Phase 2 golden dataset creation with SME reviews",
        "   ▸ Plan Layer 2 transition evaluation design",
    ])

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 24 — Thank You
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, BRAND_DARK)
    _add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
                 "Thank You", font_size=44, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
                 "Questions & Discussion",
                 font_size=24, color=BRAND_MID, alignment=PP_ALIGN.CENTER)

    # ── Save ──────────────────────────────────────────────────────
    output_path = os.path.join(os.path.dirname(__file__),
                               "wega_evals_meeting.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path


if __name__ == "__main__":
    generate_pptx()
