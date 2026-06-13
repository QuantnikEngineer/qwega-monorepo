"""Source-document ingestion + image extraction tools.

This module is the single ingestion entry point for the user-manual pipeline.
It is intentionally side-effectful: as part of reading source documents, it
also extracts and classifies every embedded image so downstream agents can
ground their content in real visual artefacts (architecture diagrams, Figma
designs, screenshots).

Responsibilities
----------------
1. Read all source documents from a single source (Confluence page or
   SharePoint folder/file). Supported types: PDF, DOCX, TXT, XLS/XLSX, JSON,
   CSV, PPTX, plus standalone image files.
2. Persist every embedded image (PDF/DOCX/PPTX) and standalone image to a
   per-session folder under ``data/extracted_images/<session_id>/``.
3. Classify each image as one of: ``architecture``, ``figma``, ``screenshot``.
   Hybrid classifier: deterministic-first, Gemini-Vision fallback when cues
   are inconclusive.
4. Push raw_corpus and image metadata into the agent session state via
   ``ToolContext`` so downstream agents read from a single source of truth.
"""

from __future__ import annotations

import io
import json
import logging
import os
import re
import tempfile
import uuid
import zipfile
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple
from urllib.parse import parse_qs, quote, unquote, urlparse

import msal
import pandas as pd
import pdfplumber
import requests
from docx import Document
from pptx import Presentation
from PIL import Image as PILImage

logger = logging.getLogger(__name__)

DATA_DIR = Path(__file__).resolve().parent.parent / "data"
IMAGES_DIR = DATA_DIR / "extracted_images"

SUPPORTED_DOC_EXT = {".pdf", ".docx", ".txt", ".xls", ".xlsx", ".json", ".csv", ".pptx"}
SUPPORTED_IMG_EXT = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".svg"}

ARCH_KEYWORDS = [
    "architecture", "arch", "topology", "deployment", "c4", "context diagram",
    "data flow", "dataflow", "block diagram", "system diagram", "components",
    "high-level design", "hld", "low-level design", "lld", "sequence diagram",
    "infra", "infrastructure", "network diagram", "er diagram", "entity relationship",
]
FIGMA_KEYWORDS = [
    "figma", "design", "wireframe", "mockup", "prototype", "ui", "ux",
    "screen", "page layout", "interaction", "user interface",
]
SCREENSHOT_KEYWORDS = [
    "screenshot", "screen capture", "snippet", "screen grab",
]

# Used to pick a landing/login image for the Getting Started section.
LANDING_KEYWORDS = [
    "login", "sign in", "sign-in", "signin", "log in", "log-in",
    "landing", "home", "homepage", "home page", "welcome",
    "dashboard", "main page", "after login",
]

LANDING_SCORE_THRESHOLD = float(os.getenv("LANDING_IMAGE_SCORE_THRESHOLD", "1.0"))

MIN_IMG_WIDTH = 120
MIN_IMG_HEIGHT = 120
CAPTION_CONTEXT_CHARS = 600
MAX_VISION_CLASSIFICATIONS = int(os.getenv("MAX_VISION_CLASSIFICATIONS", "12"))

# Cap aggregated corpus size so a 200-doc folder cannot blow Gemini's
# token limit silently. When exceeded we keep the head + tail of the
# corpus and inject an explicit truncation marker in the middle so
# downstream agents can see that some text was dropped.
MAX_CORPUS_CHARS = int(os.getenv("MAX_CORPUS_CHARS", "500000"))


def _truncate_corpus(corpus: str) -> str:
    """Return ``corpus`` if within budget, else head+marker+tail."""
    if not corpus or len(corpus) <= MAX_CORPUS_CHARS:
        return corpus
    keep_head = MAX_CORPUS_CHARS // 2
    keep_tail = MAX_CORPUS_CHARS - keep_head - 200  # leave room for marker
    head = corpus[:keep_head]
    tail = corpus[-keep_tail:] if keep_tail > 0 else ""
    marker = (
        "\n\n[...TRUNCATED... corpus exceeded "
        f"{MAX_CORPUS_CHARS} chars; "
        f"original size was {len(corpus)} chars; "
        "middle section omitted]\n\n"
    )
    logger.warning(
        "raw_corpus truncated: original=%d chars, kept=%d chars",
        len(corpus), MAX_CORPUS_CHARS,
    )
    return head + marker + tail


def _new_image_subdir() -> Path:
    IMAGES_DIR.mkdir(parents=True, exist_ok=True)
    sub = IMAGES_DIR / uuid.uuid4().hex[:12]
    sub.mkdir(parents=True, exist_ok=True)
    return sub


# ---------------------------------------------------------------------------
# SharePoint URL parsing & Graph context
# ---------------------------------------------------------------------------

def _parse_sharepoint_url(url: str) -> tuple[str, str]:
    url = url.strip().rstrip("/")
    parsed = urlparse(url)

    def _extract_from_path(raw_path: str) -> tuple[str, str]:
        raw_path = raw_path.strip("/")
        site_match = re.match(r"((?:sites|teams)/[^/]+)", raw_path, re.IGNORECASE)
        if site_match:
            site_path = "/" + site_match.group(1)
            remainder = raw_path[len(site_match.group(1)):].strip("/")
            return f"{parsed.scheme}://{parsed.netloc}{site_path}", remainder
        return "", raw_path

    qs = parse_qs(parsed.query)
    id_param = qs.get("id", [None])[0]
    if id_param:
        site_url, remainder = _extract_from_path(unquote(id_param))
        if site_url:
            return site_url, remainder

    match = re.match(r"(/(?:sites|teams)/[^/]+)", parsed.path, re.IGNORECASE)
    if match:
        site_path = match.group(1)
        remainder = parsed.path[len(site_path):].strip("/")
        if remainder.lower().endswith(".aspx"):
            remainder = ""
        site_url = f"{parsed.scheme}://{parsed.netloc}{site_path}"
        return site_url, remainder

    return url, ""


def _get_graph_context() -> dict | None:
    sharepoint_url = os.getenv("SHAREPOINT_URL", "").strip()
    if sharepoint_url:
        site_url, folder_path = _parse_sharepoint_url(sharepoint_url)
    else:
        site_url = os.getenv("SHAREPOINT_SITE_URL", "").rstrip("/")
        folder_path = os.getenv("SHAREPOINT_FOLDER_PATH", "").strip("/")
    client_id = os.getenv("SHAREPOINT_CLIENT_ID")
    client_secret = os.getenv("SHAREPOINT_CLIENT_SECRET")
    tenant_id = os.getenv("SHAREPOINT_TENANT_ID")

    if not all([site_url, folder_path, client_id, client_secret, tenant_id]):
        return None

    authority = f"https://login.microsoftonline.com/{tenant_id}"
    msal_app = msal.ConfidentialClientApplication(
        client_id=client_id, client_credential=client_secret, authority=authority,
    )
    token = msal_app.acquire_token_for_client(scopes=["https://graph.microsoft.com/.default"])
    if "access_token" not in token:
        return None

    headers = {
        "Authorization": f"Bearer {token['access_token']}",
        "Accept": "application/json",
    }

    parsed = urlparse(site_url)
    hostname = parsed.netloc
    site_resp = requests.get(
        f"https://graph.microsoft.com/v1.0/sites/{hostname}:{parsed.path}",
        headers=headers, timeout=30,
    )
    if site_resp.status_code == 401:
        raise PermissionError(
            "Microsoft Graph returned 401 Unauthorized. Ensure the Azure AD app "
            "has 'Sites.ReadWrite.All' Application permission with admin consent."
        )
    site_resp.raise_for_status()
    site_id = site_resp.json()["id"]

    path_parts = folder_path.split("/", 1)
    library_name = path_parts[0]
    path_in_drive = path_parts[1] if len(path_parts) > 1 else ""

    drives_resp = requests.get(
        f"https://graph.microsoft.com/v1.0/sites/{site_id}/drives",
        headers=headers, timeout=30,
    )
    drives_resp.raise_for_status()
    drives = drives_resp.json().get("value", [])

    drive_id = None
    for drive in drives:
        d_name = drive.get("name", "").lower()
        d_web = drive.get("webUrl", "").lower()
        if d_name == library_name.lower() or library_name.lower().replace(" ", "%20") in d_web or library_name.lower() in d_web:
            drive_id = drive["id"]
            break
    if drive_id is None and drives:
        drive_id = drives[0]["id"]
    if not drive_id:
        return None

    return {
        "headers": headers,
        "drive_id": drive_id,
        "path_in_drive": path_in_drive,
        "site_url": site_url,
        "base_folder": folder_path,
    }


# ---------------------------------------------------------------------------
# Image classification (deterministic + Gemini-Vision fallback)
# ---------------------------------------------------------------------------

def _tok(s: str) -> str:
    return (s or "").strip().lower()


def _hits(text: str, vocab: List[str]) -> int:
    text = _tok(text)
    if not text:
        return 0
    return sum(1 for k in vocab if k and (re.search(rf"\b{re.escape(k)}\b", text) or k in text))


def _classify_deterministic(filename: str, source_doc: str, caption: str) -> Optional[str]:
    haystack = " ".join(filter(None, [filename, source_doc, caption]))
    arch = _hits(haystack, ARCH_KEYWORDS)
    fig = _hits(haystack, FIGMA_KEYWORDS)
    shot = _hits(haystack, SCREENSHOT_KEYWORDS)

    scores = {"architecture": arch, "figma": fig, "screenshot": shot}
    top = max(scores.values())
    if top == 0:
        return None
    winners = [k for k, v in scores.items() if v == top]
    if len(winners) == 1 and top >= 2:
        return winners[0]
    if len(winners) == 1 and top == 1 and winners[0] in {"architecture", "figma"}:
        return winners[0]
    return None


import threading

_VISION_CALLS_THIS_RUN = {"count": 0}

# Independent budget for "describe" calls used to enrich image captions.
# Defaults to the same cap as classification but tunable separately.
MAX_VISION_DESCRIPTIONS = int(os.getenv("MAX_VISION_DESCRIPTIONS", "12"))
_VISION_DESC_CALLS = {"count": 0}

# Thread-safe access to the per-run counters when classification /
# description calls fan out across a thread pool.
_VISION_LOCK = threading.Lock()
VISION_CONCURRENCY = int(os.getenv("VISION_CONCURRENCY", "6"))


def _try_consume_budget(counter: Dict[str, int], cap: int) -> bool:
    """Atomically reserve one slot from ``counter`` if the cap allows it."""
    with _VISION_LOCK:
        if counter["count"] >= cap:
            return False
        counter["count"] += 1
        return True


def _vision_call(image_path: Path, prompt: str, counter: Dict[str, int], cap: int) -> str:
    """Single Gemini Vision call. Returns raw response text or '' on failure."""
    if not _try_consume_budget(counter, cap):
        return ""
    try:
        from google.genai import Client, types as genai_types
        client = Client()
        model = os.getenv("GEMINI_MODEL", "")
        if not model:
            return ""
        with image_path.open("rb") as fh:
            img_bytes = fh.read()
        suffix = image_path.suffix.lower().lstrip(".") or "png"
        if suffix == "svg":
            mime = "image/svg+xml"
        elif suffix == "jpg":
            mime = "image/jpeg"
        else:
            mime = f"image/{suffix}"
        response = client.models.generate_content(
            model=model,
            contents=[
                genai_types.Part.from_bytes(data=img_bytes, mime_type=mime),
                prompt,
            ],
        )
        return (getattr(response, "text", "") or "").strip()
    except Exception as exc:
        logger.warning("Gemini vision call failed: %s", exc)
        # Refund the slot so a subsequent image can use it.
        with _VISION_LOCK:
            if counter["count"] > 0:
                counter["count"] -= 1
        return ""


def _classify_with_vision(image_path: Path, caption: str) -> str:
    prompt = (
        "Classify this image into EXACTLY ONE of: architecture, figma, screenshot. "
        "Use 'architecture' for system/data-flow/deployment/topology diagrams. "
        "Use 'figma' for UI mockups/wireframes/prototypes. "
        "Use 'screenshot' for everything else. "
        f"Caption hint: {caption[:300]}\n"
        "Respond with ONE WORD only."
    )
    text = _vision_call(image_path, prompt, _VISION_CALLS_THIS_RUN, MAX_VISION_CLASSIFICATIONS).lower()
    for label in ("architecture", "figma", "screenshot"):
        if label in text:
            return label
    return "screenshot"


def _describe_with_vision(image_path: Path, existing_caption: str) -> str:
    """Return a short, role-relevance-friendly description.

    The description is intended to enrich the image's caption so the
    role-matching scorer (which works on filename + source_document +
    caption) has real semantic signal even when the surrounding document
    text is sparse (e.g. screenshots embedded in a PPTX with no nearby
    prose).
    """
    prompt = (
        "Describe this UI image in ONE short line (max 25 words). "
        "Focus on: which user role would use this screen, what action is "
        "shown (e.g. approval, dashboard, login, settings), and any visible "
        "labels or section names. Do NOT describe colours or layout. "
        "If the image is not a UI, output a single word: 'non-ui'.\n"
        f"Existing caption hint: {existing_caption[:200]}"
    )
    desc = _vision_call(image_path, prompt, _VISION_DESC_CALLS, MAX_VISION_DESCRIPTIONS)
    if desc.lower().startswith("non-ui"):
        return ""
    return desc[:300]


def _classify_image(filename: str, source_doc: str, caption: str, image_path: Path) -> str:
    label = _classify_deterministic(filename, source_doc, caption)
    if label:
        return label
    return _classify_with_vision(image_path, caption)


# ---------------------------------------------------------------------------
# Image-extraction helpers
# ---------------------------------------------------------------------------

def _save_image_bytes(out_dir: Path, base_name: str, ext: str, data: bytes) -> Optional[Path]:
    if not data:
        return None
    safe_ext = (ext or ".png").lower()
    if safe_ext not in SUPPORTED_IMG_EXT:
        safe_ext = ".png"
    out_path = out_dir / f"{base_name}{safe_ext}"
    counter = 1
    while out_path.exists():
        out_path = out_dir / f"{base_name}_{counter}{safe_ext}"
        counter += 1
    try:
        out_path.write_bytes(data)
    except OSError as exc:
        logger.warning("Failed to save image %s: %s", out_path, exc)
        return None
    if safe_ext != ".svg":
        try:
            with PILImage.open(out_path) as im:
                if im.width < MIN_IMG_WIDTH or im.height < MIN_IMG_HEIGHT:
                    out_path.unlink(missing_ok=True)
                    return None
        except Exception:
            out_path.unlink(missing_ok=True)
            return None
    return out_path


def _extract_pdf_images(pdf_path: Path, out_dir: Path, source_doc: str) -> List[Dict[str, Any]]:
    items: List[Dict[str, Any]] = []
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                page_text = page.extract_text() or ""
                for img_num, img in enumerate(page.images or []):
                    try:
                        cropped = page.crop((img["x0"], img["top"], img["x1"], img["bottom"])).to_image()
                        buf = io.BytesIO()
                        cropped.save(buf, format="PNG")
                        out_path = _save_image_bytes(
                            out_dir,
                            f"{pdf_path.stem}_p{page_num}_i{img_num}",
                            ".png",
                            buf.getvalue(),
                        )
                        if out_path is None:
                            continue
                        items.append({
                            "path": out_path.as_posix(),
                            "filename": out_path.name,
                            "source_document": source_doc,
                            "page": page_num,
                            "caption": page_text[:CAPTION_CONTEXT_CHARS],
                        })
                    except Exception as exc:
                        logger.debug("PDF image extract skipped: %s", exc)
    except Exception as exc:
        logger.warning("Could not open PDF for image extraction (%s): %s", pdf_path, exc)
    return items


def _extract_docx_images(docx_path: Path, out_dir: Path, source_doc: str) -> List[Dict[str, Any]]:
    items: List[Dict[str, Any]] = []
    try:
        with zipfile.ZipFile(docx_path) as zf:
            for member in zf.namelist():
                if not member.startswith("word/media/"):
                    continue
                ext = Path(member).suffix.lower()
                if ext not in SUPPORTED_IMG_EXT:
                    continue
                try:
                    data = zf.read(member)
                except Exception:
                    continue
                out_path = _save_image_bytes(
                    out_dir, f"{docx_path.stem}_{Path(member).stem}", ext, data,
                )
                if out_path is None:
                    continue
                items.append({
                    "path": out_path.as_posix(),
                    "filename": out_path.name,
                    "source_document": source_doc,
                    "page": None,
                    "caption": "",
                })
    except Exception as exc:
        logger.warning("Could not open DOCX for image extraction (%s): %s", docx_path, exc)
    return items


def _extract_pptx_images(pptx_path: Path, out_dir: Path, source_doc: str) -> List[Dict[str, Any]]:
    items: List[Dict[str, Any]] = []
    try:
        prs = Presentation(str(pptx_path))
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            slide_text_parts.append(paragraph.text.strip())
            slide_text = " ".join(slide_text_parts)[:CAPTION_CONTEXT_CHARS]
            for shape_num, shape in enumerate(slide.shapes):
                shape_type = getattr(shape, "shape_type", None)
                if shape_type is None:
                    continue
                # MSO_SHAPE_TYPE.PICTURE == 13
                if int(shape_type) != 13:
                    continue
                try:
                    image = shape.image
                    ext = "." + (image.ext or "png").lower()
                    out_path = _save_image_bytes(
                        out_dir, f"{pptx_path.stem}_s{slide_num}_p{shape_num}", ext, image.blob,
                    )
                    if out_path is None:
                        continue
                    items.append({
                        "path": out_path.as_posix(),
                        "filename": out_path.name,
                        "source_document": source_doc,
                        "page": slide_num,
                        "caption": slide_text,
                    })
                except Exception as exc:
                    logger.debug("PPTX image extract skipped: %s", exc)
    except Exception as exc:
        logger.warning("Could not open PPTX for image extraction (%s): %s", pptx_path, exc)
    return items


# ---------------------------------------------------------------------------
# Per-extension text readers
# ---------------------------------------------------------------------------

def _read_pdf(path: Path) -> str:
    text: List[str] = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text.append(page.extract_text() or "")
    return "\n".join(text)


def _read_docx(path: Path) -> str:
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _read_txt(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return path.read_text(errors="ignore")


def _read_excel(path: Path) -> str:
    sheets = pd.read_excel(path, sheet_name=None)
    output: List[str] = []
    for sheet_name, df in sheets.items():
        output.append(f"\n--- Sheet: {sheet_name} ---\n")
        output.append(df.to_string(index=False))
    return "\n".join(output)


def _read_json(path: Path) -> str:
    with path.open(encoding="utf-8") as fh:
        data = json.load(fh)
    return json.dumps(data, indent=2, ensure_ascii=False)


def _read_csv(path: Path) -> str:
    df = pd.read_csv(path)
    return df.to_string(index=False)


def _read_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    slides_text: List[str] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        parts = [f"\n--- Slide {slide_num} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        parts.append(text)
        slides_text.append("\n".join(parts))
    return "\n".join(slides_text)


def _parse_doc_to_text_and_images(
    local_path: Path, source_doc: str, suffix: str, images_dir: Path,
) -> Tuple[str, List[Dict[str, Any]]]:
    text = ""
    images: List[Dict[str, Any]] = []
    if suffix == ".pdf":
        text = _read_pdf(local_path)
        images = _extract_pdf_images(local_path, images_dir, source_doc)
    elif suffix == ".docx":
        text = _read_docx(local_path)
        images = _extract_docx_images(local_path, images_dir, source_doc)
    elif suffix == ".pptx":
        text = _read_pptx(local_path)
        images = _extract_pptx_images(local_path, images_dir, source_doc)
    elif suffix in {".xls", ".xlsx"}:
        text = _read_excel(local_path)
    elif suffix == ".txt":
        text = _read_txt(local_path)
    elif suffix == ".json":
        text = _read_json(local_path)
    elif suffix == ".csv":
        text = _read_csv(local_path)
    return text, images


# ---------------------------------------------------------------------------
# SharePoint reader
# ---------------------------------------------------------------------------

def _read_from_sharepoint(images_dir: Path) -> Dict[str, Any] | None:
    ctx = _get_graph_context()
    if ctx is None:
        return None

    headers = ctx["headers"]
    drive_id = ctx["drive_id"]
    path_in_drive = ctx["path_in_drive"]
    folder_path = ctx["base_folder"]

    documents: List[str] = []
    images_meta: List[Dict[str, Any]] = []

    def _download(url: str, is_preauth: bool) -> bytes:
        dl_headers = {} if is_preauth else headers
        resp = requests.get(url, headers=dl_headers, timeout=120)
        resp.raise_for_status()
        return resp.content

    def _process_doc(file_name: str, suffix: str, url: str, is_preauth: bool) -> Tuple[str, List[Dict[str, Any]]]:
        data = _download(url, is_preauth)
        tmp = Path(tempfile.mktemp(suffix=suffix))
        try:
            tmp.write_bytes(data)
            return _parse_doc_to_text_and_images(tmp, file_name, suffix, images_dir)
        finally:
            tmp.unlink(missing_ok=True)

    def _process_image(file_name: str, suffix: str, url: str, is_preauth: bool) -> Optional[Dict[str, Any]]:
        data = _download(url, is_preauth)
        out_path = _save_image_bytes(images_dir, Path(file_name).stem, suffix, data)
        if out_path is None:
            return None
        return {
            "path": out_path.as_posix(),
            "filename": out_path.name,
            "source_document": file_name,
            "page": None,
            "caption": "",
        }

    def _navigate_segments(segments: list[str]) -> dict:
        parent_url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/root/children"
        current_item: dict = {}
        for segment in segments:
            nav_resp = requests.get(parent_url, headers=headers, timeout=60)
            nav_resp.raise_for_status()
            children = nav_resp.json().get("value", [])
            matched = next(
                (c for c in children if c.get("name", "").lower() == segment.lower()), None,
            )
            if matched is None:
                raise FileNotFoundError(
                    f"Item '{segment}' not found. Available: {[c['name'] for c in children]}"
                )
            current_item = matched
            parent_url = (
                f"https://graph.microsoft.com/v1.0/drives/{drive_id}"
                f"/items/{matched['id']}/children"
            )
        return current_item

    item_suffix = Path(path_in_drive).suffix.lower() if path_in_drive else ""

    if item_suffix in SUPPORTED_DOC_EXT or item_suffix in SUPPORTED_IMG_EXT:
        encoded_path = quote(path_in_drive, safe="/")
        item_resp = requests.get(
            f"https://graph.microsoft.com/v1.0/drives/{drive_id}/root:/{encoded_path}",
            headers=headers, timeout=60,
        )
        if item_resp.status_code == 404:
            item_data = _navigate_segments(path_in_drive.split("/"))
        else:
            item_resp.raise_for_status()
            item_data = item_resp.json()

        file_name = item_data.get("name", Path(path_in_drive).name)
        download_url = item_data.get("@microsoft.graph.downloadUrl")
        is_preauth = bool(download_url)
        if not download_url:
            download_url = (
                f"https://graph.microsoft.com/v1.0/drives/{drive_id}"
                f"/items/{item_data['id']}/content"
            )
        if item_suffix in SUPPORTED_DOC_EXT:
            text, imgs = _process_doc(file_name, item_suffix, download_url, is_preauth)
            documents.append(text)
            images_meta.extend(imgs)
        else:
            meta = _process_image(file_name, item_suffix, download_url, is_preauth)
            if meta:
                images_meta.append(meta)
        return {
            "status": "success",
            "documents": documents,
            "images": images_meta,
            "source": "sharepoint",
            "folder": folder_path,
        }

    if path_in_drive:
        encoded_path = quote(path_in_drive, safe="/")
        list_url = (
            f"https://graph.microsoft.com/v1.0/drives/{drive_id}"
            f"/root:/{encoded_path}:/children"
        )
    else:
        list_url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/root/children"

    list_resp = requests.get(list_url, headers=headers, timeout=60)
    if list_resp.status_code == 404 and path_in_drive:
        folder_segments = path_in_drive.split("/")
        parent_url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/root/children"
        current_id = None
        for segment in folder_segments:
            nav_resp = requests.get(parent_url, headers=headers, timeout=60)
            nav_resp.raise_for_status()
            children = nav_resp.json().get("value", [])
            matched = next(
                (c for c in children if c.get("name", "").lower() == segment.lower() and "folder" in c),
                None,
            )
            if matched is None:
                raise FileNotFoundError(
                    f"Folder segment '{segment}' not found. Available: {[c['name'] for c in children]}"
                )
            current_id = matched["id"]
            parent_url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/items/{current_id}/children"
        list_resp = requests.get(parent_url, headers=headers, timeout=60)

    list_resp.raise_for_status()
    items = list_resp.json().get("value", [])

    docs_to_process: List[Tuple[str, str, str, bool]] = []
    images_to_process: List[Tuple[str, str, str, bool]] = []
    for item in items:
        if "folder" in item:
            continue
        file_name = item.get("name", "")
        suffix = Path(file_name).suffix.lower()
        download_url = item.get("@microsoft.graph.downloadUrl")
        is_preauth = bool(download_url)
        if not download_url:
            item_id = item["id"]
            download_url = (
                f"https://graph.microsoft.com/v1.0/drives/{drive_id}"
                f"/items/{item_id}/content"
            )
        if suffix in SUPPORTED_DOC_EXT:
            docs_to_process.append((file_name, suffix, download_url, is_preauth))
        elif suffix in SUPPORTED_IMG_EXT:
            images_to_process.append((file_name, suffix, download_url, is_preauth))

    with ThreadPoolExecutor(max_workers=min(8, len(docs_to_process) or 1)) as pool:
        futures = {
            pool.submit(_process_doc, name, sfx, url, preauth): name
            for name, sfx, url, preauth in docs_to_process
        }
        for future in as_completed(futures):
            try:
                text, imgs = future.result()
                documents.append(text)
                images_meta.extend(imgs)
            except Exception as exc:
                logger.warning("Failed to process SharePoint doc: %s", exc)

    for name, sfx, url, preauth in images_to_process:
        try:
            meta = _process_image(name, sfx, url, preauth)
            if meta:
                images_meta.append(meta)
        except Exception as exc:
            logger.warning("Failed to download SharePoint image %s: %s", name, exc)

    return {
        "status": "success",
        "documents": documents,
        "images": images_meta,
        "source": "sharepoint",
        "folder": folder_path,
    }


# ---------------------------------------------------------------------------
# Confluence reader (text + attachments)
# ---------------------------------------------------------------------------

def _confluence_auth() -> Tuple[Dict[str, str], str, str]:
    email = os.getenv("CONFLUENCE_EMAIL", "").strip()
    api_token = os.getenv("CONFLUENCE_API_TOKEN", "").strip()
    if not email or not api_token:
        raise RuntimeError(
            "CONFLUENCE_EMAIL and CONFLUENCE_API_TOKEN must be set to read Confluence pages."
        )
    import base64 as _b64
    raw = f"{email}:{api_token}".encode("utf-8")
    encoded = _b64.b64encode(raw).decode("ascii")
    return {"Authorization": f"Basic {encoded}", "Accept": "application/json"}, email, api_token


def _parse_confluence_page_url(url: str) -> tuple[str, str]:
    parsed = urlparse(url)
    base_url = f"{parsed.scheme}://{parsed.netloc}"
    match = re.search(r"/pages/(\d+)", parsed.path)
    return base_url, (match.group(1) if match else "")


def _html_to_text(html: str) -> str:
    from html.parser import HTMLParser

    class _T(HTMLParser):
        BLOCK = {"p", "br", "h1", "h2", "h3", "h4", "h5", "h6", "li", "tr", "div", "section", "article", "pre"}
        SKIP = {"script", "style"}

        def __init__(self):
            super().__init__()
            self.parts: List[str] = []
            self._skip = False

        def handle_starttag(self, tag, attrs):
            if tag in self.SKIP:
                self._skip = True
            elif tag in self.BLOCK:
                self.parts.append("\n")

        def handle_endtag(self, tag):
            if tag in self.SKIP:
                self._skip = False

        def handle_data(self, data):
            if not self._skip:
                self.parts.append(data)

        def text(self) -> str:
            return re.sub(r"\n{3,}", "\n\n", "".join(self.parts)).strip()

    p = _T()
    p.feed(html or "")
    return p.text()


def _read_confluence_attachments(
    base_url: str, page_id: str, auth_headers: Dict[str, str],
    page_title: str, images_dir: Path,
) -> List[Dict[str, Any]]:
    items: List[Dict[str, Any]] = []
    endpoint = f"{base_url}/wiki/rest/api/content/{page_id}/child/attachment?limit=200"
    try:
        resp = requests.get(endpoint, headers=auth_headers, timeout=30)
        if not resp.ok:
            return []
        for att in resp.json().get("results", []):
            title = att.get("title", "") or ""
            ext = Path(title).suffix.lower()
            if ext not in SUPPORTED_IMG_EXT:
                continue
            download_rel = (att.get("_links") or {}).get("download", "")
            if not download_rel:
                continue
            download_url = f"{base_url}/wiki{download_rel}" if download_rel.startswith("/") else download_rel
            try:
                d_resp = requests.get(download_url, headers=auth_headers, timeout=60)
                d_resp.raise_for_status()
            except Exception:
                continue
            out_path = _save_image_bytes(images_dir, Path(title).stem, ext, d_resp.content)
            if out_path is None:
                continue
            items.append({
                "path": out_path.as_posix(),
                "filename": out_path.name,
                "source_document": f"Confluence: {page_title}",
                "page": None,
                "caption": title,
            })
    except Exception as exc:
        logger.warning("Failed to read Confluence attachments for %s: %s", page_id, exc)
    return items


def _read_from_confluence(url: str, images_dir: Path) -> Dict[str, Any]:
    base_url, page_id = _parse_confluence_page_url(url)
    if not page_id:
        raise ValueError(f"Could not extract page ID from Confluence URL: {url}")
    headers, _, _ = _confluence_auth()

    def _fetch_page(pid: str) -> Tuple[str, str]:
        endpoint = f"{base_url}/wiki/rest/api/content/{pid}?expand=body.export_view,title"
        resp = requests.get(endpoint, headers=headers, timeout=30)
        if resp.status_code >= 400:
            raise RuntimeError(
                f"Confluence API returned {resp.status_code} for page {pid}: {resp.text}"
            )
        data = resp.json()
        title = data.get("title", "")
        html_body = (data.get("body") or {}).get("export_view", {}).get("value", "")
        text = _html_to_text(html_body)
        # Do NOT prepend the Confluence page title as a markdown heading.
        # The Confluence page title is typically a filing label
        # (e.g. the previous run's project_name) and would bias the
        # extraction LLM into echoing it as document_title. We still
        # return the title alongside for attachment bookkeeping, but the
        # corpus we feed to extraction contains only the body content.
        return title, text

    documents: List[str] = []
    images_meta: List[Dict[str, Any]] = []

    title, text = _fetch_page(page_id)
    documents.append(text)
    images_meta.extend(_read_confluence_attachments(base_url, page_id, headers, title, images_dir))

    children_url = f"{base_url}/wiki/rest/api/content/{page_id}/child/page?limit=50&expand=title"
    try:
        children_resp = requests.get(children_url, headers=headers, timeout=30)
        if children_resp.ok:
            for child in children_resp.json().get("results", []):
                cid = child.get("id")
                if not cid:
                    continue
                try:
                    c_title, c_text = _fetch_page(cid)
                    documents.append(c_text)
                    images_meta.extend(_read_confluence_attachments(base_url, cid, headers, c_title, images_dir))
                except Exception as exc:
                    logger.warning("Skipping inaccessible Confluence child %s: %s", cid, exc)
    except Exception as exc:
        logger.warning("Failed to enumerate Confluence children for %s: %s", page_id, exc)

    return {
        "status": "success",
        "documents": documents,
        "images": images_meta,
        "source": "confluence_input",
        "folder": url,
    }


# Public alias retained for backwards compat
def read_confluence_page_content(url: str) -> Dict[str, Any]:
    images_dir = _new_image_subdir()
    return _read_from_confluence(url, images_dir)


# ---------------------------------------------------------------------------
# Local data reader (dev fallback)
# ---------------------------------------------------------------------------

def _read_from_local(images_dir: Path) -> Dict[str, Any]:
    documents: List[str] = []
    images_meta: List[Dict[str, Any]] = []
    if not DATA_DIR.exists():
        return {"status": "success", "documents": documents, "images": images_meta,
                "source": "local", "folder": str(DATA_DIR)}
    for path in DATA_DIR.iterdir():
        if path.is_dir():
            continue
        suffix = path.suffix.lower()
        if suffix in SUPPORTED_DOC_EXT:
            try:
                text, imgs = _parse_doc_to_text_and_images(path, path.name, suffix, images_dir)
                documents.append(text)
                images_meta.extend(imgs)
            except Exception as exc:
                logger.warning("Failed to process local doc %s: %s", path, exc)
        elif suffix in SUPPORTED_IMG_EXT:
            try:
                data = path.read_bytes()
                out_path = _save_image_bytes(images_dir, path.stem, suffix, data)
                if out_path:
                    images_meta.append({
                        "path": out_path.as_posix(),
                        "filename": out_path.name,
                        "source_document": path.name,
                        "page": None,
                        "caption": "",
                    })
            except Exception as exc:
                logger.warning("Failed to read local image %s: %s", path, exc)
    return {
        "status": "success",
        "documents": documents,
        "images": images_meta,
        "source": "local",
        "folder": str(DATA_DIR),
    }


# ---------------------------------------------------------------------------
# Public ingestion tool
# ---------------------------------------------------------------------------

def _select_getting_started_image(images: List[Dict[str, Any]]) -> Optional[Dict[str, Any]]:
    """Pick the best login/landing/home image for the Getting Started section.

    Considers only figma/screenshot images (architecture diagrams are never
    placed in Getting Started). Scores filename + source_document + caption
    + description against ``LANDING_KEYWORDS``; returns the top-scoring
    image if its score clears ``LANDING_SCORE_THRESHOLD``, else None.
    """
    candidates = [i for i in images if i.get("category") in {"figma", "screenshot"}]
    if not candidates:
        return None

    scored: List[Tuple[float, Dict[str, Any]]] = []
    for img in candidates:
        fname = (img.get("filename") or "").lower()
        src_doc = (img.get("source_document") or "").lower()
        caption = (img.get("caption") or "").lower()
        description = (img.get("description") or "").lower()

        score = 0.0
        score += 1.2 * _hits(fname, LANDING_KEYWORDS)
        score += 0.8 * _hits(src_doc, LANDING_KEYWORDS)
        score += 0.7 * _hits(caption, LANDING_KEYWORDS)
        score += 1.4 * _hits(description, LANDING_KEYWORDS)
        if img.get("category") == "figma":
            score += 0.25

        if score > 0:
            scored.append((score, img))

    if not scored:
        return None
    scored.sort(key=lambda x: x[0], reverse=True)
    top_score, top_img = scored[0]
    if top_score < LANDING_SCORE_THRESHOLD:
        return None
    return {
        "path": top_img.get("path"),
        "filename": top_img.get("filename"),
        "source_document": top_img.get("source_document"),
        "caption": top_img.get("caption"),
        "description": top_img.get("description"),
        "category": top_img.get("category"),
        "score": round(top_score, 2),
    }


_CAPTION_DESCRIPTION_THRESHOLD = int(os.getenv("CAPTION_DESCRIPTION_THRESHOLD", "120"))


def _process_one_image(img: Dict[str, Any]) -> Dict[str, Any]:
    """Classify + (optionally) caption-enrich a single image.

    Pure per-image work so we can fan it out across a ThreadPoolExecutor.
    Mutates and returns ``img``.
    """
    path = img.get("path")
    if not path:
        return img
    existing_caption = img.get("caption", "") or ""
    label = _classify_image(
        filename=img.get("filename", ""),
        source_doc=img.get("source_document", ""),
        caption=existing_caption,
        image_path=Path(path),
    )
    img["category"] = label

    # Only enrich figma/screenshot images whose caption is too thin to
    # ground role matching. CAPTION_DESCRIPTION_THRESHOLD is tunable.
    if label in {"figma", "screenshot"} and len(existing_caption) < _CAPTION_DESCRIPTION_THRESHOLD:
        desc = _describe_with_vision(Path(path), existing_caption)
        if desc:
            img["description"] = desc
            img["caption"] = (existing_caption + " " + desc).strip()
    return img


def _classify_all(images: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    """Classify every image and (for UI-class images with weak captions)
    fetch a short Vision-generated description, merging it into the caption.

    Vision calls run concurrently across a ThreadPoolExecutor to cut total
    pipeline time. Per-run budgets (MAX_VISION_CLASSIFICATIONS,
    MAX_VISION_DESCRIPTIONS) are still respected via thread-safe counters.
    """
    deduped: List[Dict[str, Any]] = []
    seen = set()
    for img in images:
        path = img.get("path")
        if not path or path in seen:
            continue
        seen.add(path)
        deduped.append(img)

    if not deduped:
        return []

    workers = max(1, min(VISION_CONCURRENCY, len(deduped)))
    results: List[Dict[str, Any]] = [None] * len(deduped)  # type: ignore[list-item]
    with ThreadPoolExecutor(max_workers=workers) as pool:
        future_to_idx = {pool.submit(_process_one_image, img): i for i, img in enumerate(deduped)}
        for future in as_completed(future_to_idx):
            idx = future_to_idx[future]
            try:
                results[idx] = future.result()
            except Exception as exc:
                logger.warning("Image processing failed: %s", exc)
                results[idx] = deduped[idx]
    return [r for r in results if r is not None]


def read_all_files_from_data(tool_context: Any = None) -> Dict[str, Any]:
    """Read all source documents AND extract embedded images.

    When invoked by an LlmAgent, ``tool_context.state`` is populated with:
    ``raw_corpus``, ``images``, ``architecture_images``, ``figma_images``,
    ``screenshot_images``, ``has_architecture_diagrams``, ``has_figma_designs``,
    ``ingestion_summary``, ``images_dir``.
    """
    images_dir = _new_image_subdir()

    confluence_source_url = os.getenv("CONFLUENCE_SOURCE_URL", "").strip()
    if confluence_source_url:
        result = _read_from_confluence(confluence_source_url, images_dir)
    else:
        result = _read_from_sharepoint(images_dir)
        if result is None:
            result = _read_from_local(images_dir)
            if not result["documents"]:
                raise RuntimeError(
                    "No source configured. Set CONFLUENCE_SOURCE_URL for a Confluence page, "
                    "or SHAREPOINT_URL (with SHAREPOINT_CLIENT_ID, SHAREPOINT_CLIENT_SECRET, "
                    "SHAREPOINT_TENANT_ID) for a SharePoint folder/file."
                )

    documents: List[str] = result.get("documents", [])
    images: List[Dict[str, Any]] = _classify_all(result.get("images", []))
    raw_corpus = "\n\n".join(d for d in documents if d).strip()
    raw_corpus = _truncate_corpus(raw_corpus)

    arch_imgs = [i for i in images if i.get("category") == "architecture"]
    fig_imgs = [i for i in images if i.get("category") == "figma"]
    shot_imgs = [i for i in images if i.get("category") == "screenshot"]
    getting_started_image = _select_getting_started_image(images)

    summary = {
        "source": result.get("source"),
        "documents_count": len(documents),
        "images_total": len(images),
        "architecture_images": len(arch_imgs),
        "figma_images": len(fig_imgs),
        "screenshot_images": len(shot_imgs),
        "getting_started_image": bool(getting_started_image),
        "images_dir": images_dir.as_posix(),
    }
    logger.info("Ingestion summary: %s", summary)

    state_obj = getattr(tool_context, "state", None)
    if state_obj is not None:
        state_obj["raw_corpus"] = raw_corpus
        state_obj["images"] = images
        state_obj["architecture_images"] = arch_imgs
        state_obj["figma_images"] = fig_imgs
        state_obj["screenshot_images"] = shot_imgs
        state_obj["getting_started_image"] = getting_started_image  # may be None
        state_obj["has_architecture_diagrams"] = bool(arch_imgs)
        state_obj["has_figma_designs"] = bool(fig_imgs) or bool(shot_imgs)
        state_obj["has_getting_started_image"] = bool(getting_started_image)
        state_obj["ingestion_summary"] = summary
        state_obj["images_dir"] = images_dir.as_posix()

    return {
        "status": "success",
        "raw_corpus": raw_corpus,
        "documents": documents,
        "images": images,
        "architecture_images": arch_imgs,
        "figma_images": fig_imgs,
        "screenshot_images": shot_imgs,
        "getting_started_image": getting_started_image,
        "has_architecture_diagrams": bool(arch_imgs),
        "has_figma_designs": bool(fig_imgs) or bool(shot_imgs),
        "has_getting_started_image": bool(getting_started_image),
        "summary": summary,
        "source": result.get("source"),
        "folder": result.get("folder"),
    }


# ---------------------------------------------------------------------------
# SharePoint upload (used by the optional PDF exporter pipeline)
# ---------------------------------------------------------------------------

def upload_pdf_to_sharepoint(local_pdf_path: str) -> str:
    ctx = _get_graph_context()
    if ctx is None:
        raise RuntimeError(
            "Cannot upload to SharePoint: SHAREPOINT_URL (or SHAREPOINT_SITE_URL), "
            "SHAREPOINT_CLIENT_ID, SHAREPOINT_CLIENT_SECRET, and "
            "SHAREPOINT_TENANT_ID must all be set."
        )

    headers = ctx["headers"]
    drive_id = ctx["drive_id"]
    path_in_drive = ctx["path_in_drive"]
    site_url = ctx["site_url"]
    base_folder = ctx["base_folder"]

    upload_folder = f"{path_in_drive}/user-manual" if path_in_drive else "user-manual"
    encoded_folder = quote(upload_folder, safe="/")
    folder_url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/root:/{encoded_folder}"
    folder_resp = requests.get(folder_url, headers=headers, timeout=30)
    if folder_resp.status_code == 404:
        parent_path = "/".join(upload_folder.split("/")[:-1])
        folder_name = upload_folder.split("/")[-1]
        if parent_path:
            create_url = (
                f"https://graph.microsoft.com/v1.0/drives/{drive_id}"
                f"/root:/{quote(parent_path, safe='/')}:/children"
            )
        else:
            create_url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/root/children"
        create_resp = requests.post(
            create_url,
            headers={**headers, "Content-Type": "application/json"},
            json={"name": folder_name, "folder": {}, "@microsoft.graph.conflictBehavior": "fail"},
            timeout=30,
        )
        if create_resp.status_code not in (200, 201, 409):
            create_resp.raise_for_status()

    pdf_path = Path(local_pdf_path)
    file_name = pdf_path.name
    file_size = pdf_path.stat().st_size
    encoded_upload_path = quote(f"{upload_folder}/{file_name}", safe="/")

    if file_size < 4 * 1024 * 1024:
        upload_url = (
            f"https://graph.microsoft.com/v1.0/drives/{drive_id}"
            f"/root:/{encoded_upload_path}:/content"
        )
        with pdf_path.open("rb") as fh:
            upload_resp = requests.put(
                upload_url,
                headers={**headers, "Content-Type": "application/pdf"},
                data=fh, timeout=120,
            )
        upload_resp.raise_for_status()
        web_url = upload_resp.json().get("webUrl", "")
    else:
        session_url = (
            f"https://graph.microsoft.com/v1.0/drives/{drive_id}"
            f"/root:/{encoded_upload_path}:/createUploadSession"
        )
        session_resp = requests.post(
            session_url,
            headers={**headers, "Content-Type": "application/json"},
            json={"item": {"@microsoft.graph.conflictBehavior": "replace"}},
            timeout=30,
        )
        session_resp.raise_for_status()
        upload_url = session_resp.json()["uploadUrl"]
        chunk_size = 10 * 1024 * 1024
        with pdf_path.open("rb") as fh:
            start = 0
            while start < file_size:
                chunk = fh.read(chunk_size)
                end = start + len(chunk) - 1
                chunk_headers = {
                    "Content-Length": str(len(chunk)),
                    "Content-Range": f"bytes {start}-{end}/{file_size}",
                }
                chunk_resp = requests.put(upload_url, headers=chunk_headers, data=chunk, timeout=120)
                chunk_resp.raise_for_status()
                start = end + 1
        web_url = chunk_resp.json().get("webUrl", "")

    parsed = urlparse(site_url)
    site_path = parsed.path.rstrip("/")
    server_relative = f"{site_path}/{base_folder}/user-manual/{file_name}"
    return web_url or server_relative
